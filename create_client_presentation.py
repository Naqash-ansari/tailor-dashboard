import sys
from pathlib import Path

sys.path.insert(0, str(Path(".tools/python-pptx").resolve()))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("Aans-Fabric-Tailor-Dashboard-Presentation.pptx")
LOGO = Path("shop-dashboard-test/public/brand/aans-fabric-logo-icon.png")

BRAND_DARK = RGBColor(18, 43, 42)
BRAND_GREEN = RGBColor(13, 107, 95)
BRAND_GOLD = RGBColor(216, 176, 91)
BRAND_IVORY = RGBColor(247, 244, 238)
TEXT = RGBColor(23, 32, 51)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color=BRAND_IVORY):
    shape = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill=WHITE, line=RGBColor(225, 214, 196)):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_button(slide, text, x, y, w, fill=BRAND_GOLD, color=BRAND_DARK):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color
    return shape


def add_header(slide, title, subtitle=None):
    add_text(slide, title, 0.75, 0.55, 8.5, 0.5, 24, True, BRAND_DARK)
    if subtitle:
        add_text(slide, subtitle, 0.75, 1.05, 8.6, 0.35, 11, False, MUTED)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.25), Inches(0.35), height=Inches(0.75))


def add_bullets(slide, items, x, y, w, h, size=16, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BRAND_DARK)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.85), Inches(0.7), height=Inches(1.1))
    add_text(slide, "Aans Fabric", 0.85, 2.0, 8.5, 0.6, 34, True, WHITE)
    add_text(slide, "Tailor Dashboard System", 0.85, 2.65, 8.5, 0.45, 22, True, BRAND_GOLD)
    add_text(slide, "Client presentation: proposed screens, workflow, and local data setup", 0.85, 3.25, 8.8, 0.45, 15, False, RGBColor(232, 223, 210))
    add_button(slide, "Browser based - No Electron", 0.85, 4.15, 2.4, BRAND_GOLD)
    add_button(slide, "Local SQLite data", 3.45, 4.15, 2.0, RGBColor(255, 248, 231), BRAND_DARK)
    add_button(slide, "Desktop icon launch", 5.65, 4.15, 2.15, RGBColor(255, 248, 231), BRAND_DARK)
    add_text(slide, "Prepared for client walkthrough", 0.85, 6.65, 5.0, 0.35, 12, False, RGBColor(232, 223, 210))


def overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Project Overview", "A professional local dashboard for tailor customer and measurement management.")
    cards = [
        ("Customer Records", "Save name, phone, address, order dates, notes."),
        ("Measurements", "Qameez, suit, shalwar, trouser, cm/inch support."),
        ("Payments", "Stitching price, advance, remaining balance."),
        ("Local Desktop Use", "Click desktop icon, app opens in browser at localhost:3000."),
    ]
    for idx, (title, body) in enumerate(cards):
        x = 0.75 + (idx % 2) * 6.05
        y = 1.8 + (idx // 2) * 2.0
        add_card(slide, x, y, 5.55, 1.55)
        add_text(slide, title, x + 0.25, y + 0.25, 4.6, 0.3, 17, True, BRAND_DARK)
        add_text(slide, body, x + 0.25, y + 0.72, 4.7, 0.5, 12, False, MUTED)


def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "User Flow", "How the client will use the system from desktop icon to saved record.")
    steps = [
        ("1", "Click Desktop Icon"),
        ("2", "Browser Opens"),
        ("3", "Dashboard View"),
        ("4", "Add Customer"),
        ("5", "Enter Measurements"),
        ("6", "Save to SQLite DB"),
        ("7", "Preview / Edit / Delete"),
    ]
    x = 0.7
    for number, label in steps:
        add_card(slide, x, 2.45, 1.45, 1.35, WHITE)
        add_text(slide, number, x + 0.47, 2.58, 0.5, 0.35, 20, True, BRAND_GREEN, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.12, 3.02, 1.2, 0.45, 10, True, TEXT, PP_ALIGN.CENTER)
        if number != "7":
            add_text(slide, ">", x + 1.55, 2.86, 0.3, 0.3, 18, True, BRAND_GOLD)
        x += 1.75
    add_text(slide, "Daily operation is simple: the user clicks the icon, enters orders, and the data remains saved locally.", 1.3, 5.2, 10.6, 0.5, 16, False, BRAND_DARK, PP_ALIGN.CENTER)


def dashboard_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Screen 1: Dashboard", "Aans Fabric branded dashboard with quick actions and business summary.")
    add_card(slide, 0.75, 1.55, 11.85, 5.25, WHITE)
    add_card(slide, 1.05, 1.85, 11.25, 1.25, BRAND_DARK, BRAND_DARK)
    add_text(slide, "Aans Fabric Tailor Dashboard", 1.35, 2.1, 5.5, 0.35, 20, True, WHITE)
    add_button(slide, "Add customer", 9.15, 2.15, 1.45, BRAND_GOLD)
    add_button(slide, "Saved customers", 10.75, 2.15, 1.25, RGBColor(255, 248, 231), BRAND_DARK)
    metrics = [("Saved Customers", "24"), ("Today Orders", "6"), ("Pending Payments", "Rs 18,500"), ("Delivered", "14")]
    for idx, (label, value) in enumerate(metrics):
        x = 1.05 + idx * 2.82
        add_card(slide, x, 3.35, 2.55, 1.05)
        add_text(slide, label, x + 0.18, 3.55, 2.1, 0.25, 10, True, MUTED)
        add_text(slide, value, x + 0.18, 3.9, 2.1, 0.35, 18, True, BRAND_DARK)
    add_card(slide, 1.05, 4.75, 5.4, 1.55)
    add_text(slide, "Animated Weekly Orders Chart", 1.3, 4.95, 4.5, 0.3, 15, True, BRAND_DARK)
    for idx, h in enumerate([0.45, 0.75, 0.58, 1.0, 0.86, 0.62]):
        bar = slide.shapes.add_shape(1, Inches(1.45 + idx * 0.55), Inches(6.05 - h), Inches(0.28), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BRAND_GOLD
        bar.line.fill.background()
    add_card(slide, 6.75, 4.75, 5.55, 1.55)
    add_text(slide, "Order Pipeline + Payment Focus", 7.0, 4.95, 4.9, 0.35, 15, True, BRAND_DARK)
    add_bullets(slide, ["Pending / In Progress / Completed / Delivered", "Remaining payment summary", "Recent customer records"], 7.0, 5.4, 4.9, 0.8, 10, MUTED)


def form_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Screen 2: Customer Measurement Form", "Professional entry page for customer details, measurements, and payments.")
    sections = [
        ("Customer Details", ["Name", "Phone", "Address", "Order date", "Delivery date"]),
        ("Suit / Qameez", ["Length", "Chest", "Shoulder", "Sleeve", "Collar", "Pockets"]),
        ("Shalwar / Trouser", ["Length", "Waist", "Hip", "Thigh", "Poncha", "Zip/Elastic"]),
        ("Order Details", ["Fabric type", "Suit design", "Price", "Advance", "Status"]),
    ]
    for idx, (title, fields) in enumerate(sections):
        x = 0.75 + (idx % 2) * 6.05
        y = 1.65 + (idx // 2) * 2.15
        add_card(slide, x, y, 5.55, 1.7)
        add_text(slide, title, x + 0.25, y + 0.2, 4.6, 0.3, 16, True, BRAND_DARK)
        add_bullets(slide, [" / ".join(fields[:3]), " / ".join(fields[3:])], x + 0.25, y + 0.68, 4.8, 0.7, 11, MUTED)
    add_button(slide, "CM / Inch selector", 4.3, 6.25, 1.8, BRAND_GREEN, WHITE)
    add_button(slide, "Numeric only validation", 6.35, 6.25, 2.0, BRAND_GOLD, BRAND_DARK)


def preview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Screen 3: Measurement Preview", "Live visual measurement sheet for client confidence before saving.")
    add_card(slide, 0.8, 1.55, 7.1, 5.35)
    add_text(slide, "Custom Measurement Diagram", 1.1, 1.85, 4.5, 0.35, 18, True, BRAND_DARK)
    # Simple custom mock diagram
    for cx in [2.55, 5.65]:
        head = slide.shapes.add_shape(9, Inches(cx), Inches(2.25), Inches(0.55), Inches(0.55))
        head.fill.background()
        head.line.color.rgb = BRAND_DARK
        head.line.width = Pt(2)
    for x in [2.15, 5.25]:
        body = slide.shapes.add_shape(4, Inches(x), Inches(2.95), Inches(1.35), Inches(2.05))
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(255, 250, 240)
        body.line.color.rgb = BRAND_DARK
        body.line.width = Pt(2)
    add_text(slide, "Qameez Front", 2.0, 5.25, 1.6, 0.3, 11, True, BRAND_GREEN, PP_ALIGN.CENTER)
    add_text(slide, "Shalwar / Back", 5.0, 5.25, 1.8, 0.3, 11, True, BRAND_GREEN, PP_ALIGN.CENTER)
    for idx, label in enumerate(["Length", "Chest", "Shoulder", "Sleeve", "Waist", "Poncha"]):
        x = 8.25 + (idx % 2) * 2.0
        y = 1.85 + (idx // 2) * 1.15
        add_card(slide, x, y, 1.7, 0.8)
        add_text(slide, label, x + 0.15, y + 0.15, 1.2, 0.2, 10, True, MUTED)
        add_text(slide, "Value + unit", x + 0.15, y + 0.42, 1.35, 0.2, 10, True, BRAND_DARK)
    add_bullets(slide, ["Preview appears on form", "Preview button available on saved records", "Uses our own custom art, not copied image"], 8.25, 5.45, 4.0, 1.0, 12, TEXT)


def records_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Screen 4: Saved Customers", "Searchable customer list with edit, preview, and safe delete.")
    add_card(slide, 0.85, 1.55, 11.6, 4.9)
    add_text(slide, "Saved Customers", 1.15, 1.85, 3.0, 0.35, 20, True, BRAND_DARK)
    add_card(slide, 8.3, 1.8, 3.6, 0.45)
    add_text(slide, "Search by name or phone", 8.55, 1.92, 2.5, 0.2, 10, False, MUTED)
    headers = ["Customer", "Order", "Measurements", "Payment", "Actions"]
    xs = [1.15, 3.45, 5.55, 8.0, 10.0]
    for x, h in zip(xs, headers):
        add_text(slide, h, x, 2.65, 1.6, 0.25, 10, True, MUTED)
    for row in range(3):
        y = 3.1 + row * 0.9
        add_text(slide, "Customer name\n0300 0000000", 1.15, y, 1.8, 0.5, 10, False, TEXT)
        add_text(slide, "Pending", 3.45, y, 1.2, 0.3, 10, True, BRAND_GREEN)
        add_text(slide, "Chest / Shoulder\nSleeve / Length", 5.55, y, 1.8, 0.45, 10, False, TEXT)
        add_text(slide, "Rs 3,000\nRemaining 500", 8.0, y, 1.4, 0.45, 10, False, TEXT)
        add_button(slide, "Edit", 10.0, y, 0.6, WHITE, TEXT)
        add_button(slide, "Preview", 10.7, y, 0.85, RGBColor(255, 248, 231), BRAND_DARK)
        add_button(slide, "Delete", 11.65, y, 0.7, RGBColor(255, 235, 235), RGBColor(153, 27, 27))


def data_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Data Storage: Local SQLite", "Reliable local database instead of browser-only storage.")
    add_card(slide, 0.9, 1.7, 5.7, 4.5)
    add_text(slide, "Where data is saved", 1.2, 2.0, 3.6, 0.35, 18, True, BRAND_DARK)
    add_text(slide, "shop-dashboard-test\\data\\tailor-dashboard.db", 1.2, 2.65, 4.9, 0.4, 14, True, BRAND_GREEN)
    add_bullets(slide, ["Works offline on the client laptop", "Browser clear karne se data delete nahi hota", "Backup is simple: copy the .db file", "No monthly database fee for local use"], 1.2, 3.3, 4.9, 1.8, 13, TEXT)
    add_card(slide, 7.0, 1.7, 5.2, 4.5, BRAND_DARK, BRAND_DARK)
    add_text(slide, "Client Setup", 7.35, 2.0, 3.0, 0.35, 18, True, WHITE)
    add_bullets(slide, ["One-time install on laptop", "Desktop shortcut created", "App opens in browser", "Terminal can run hidden in background"], 7.35, 2.65, 4.2, 1.5, 14, RGBColor(232, 223, 210))
    add_button(slide, "Daily use: click icon", 7.35, 5.2, 2.4, BRAND_GOLD, BRAND_DARK)


def roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Delivery Plan", "Recommended phased delivery for client approval.")
    phases = [
        ("Phase 1", "Local dashboard, form, records, SQLite, desktop icon."),
        ("Phase 2", "Print measurement slip, invoice, backup/restore button."),
        ("Phase 3", "Reports, order reminders, customer SMS/WhatsApp options."),
        ("Phase 4", "Optional online sync or multi-device version."),
    ]
    for idx, (phase, desc) in enumerate(phases):
        y = 1.65 + idx * 1.15
        add_card(slide, 1.0, y, 11.3, 0.85)
        add_text(slide, phase, 1.3, y + 0.2, 1.5, 0.25, 15, True, BRAND_GREEN)
        add_text(slide, desc, 2.9, y + 0.2, 8.4, 0.25, 14, False, TEXT)
    add_text(slide, "Goal: make order entry fast, professional, and safe for daily tailor shop use.", 1.0, 6.65, 10.8, 0.35, 15, True, BRAND_DARK, PP_ALIGN.CENTER)


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BRAND_DARK)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(5.75), Inches(0.85), height=Inches(1.15))
    add_text(slide, "Aans Fabric Tailor Dashboard", 1.4, 2.4, 10.5, 0.55, 30, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, "A clean local system for measurements, orders, payments, and customer records.", 2.05, 3.15, 9.2, 0.4, 16, False, RGBColor(232, 223, 210), PP_ALIGN.CENTER)
    add_button(slide, "Ready for client demo", 5.25, 4.15, 2.8, BRAND_GOLD, BRAND_DARK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    overview_slide(prs)
    flow_slide(prs)
    dashboard_slide(prs)
    form_slide(prs)
    preview_slide(prs)
    records_slide(prs)
    data_slide(prs)
    roadmap_slide(prs)
    closing_slide(prs)

    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
